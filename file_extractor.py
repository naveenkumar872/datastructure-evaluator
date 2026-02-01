"""
File content extraction utility for various formats.
Supports: TXT, PDF, DOC, DOCX, PPT, PPTX
"""

import os
import re


def extract_text_from_txt(file_path):
    """Extract text from TXT file"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        with open(file_path, 'r', encoding='latin-1') as f:
            return f.read()


def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    try:
        import PyPDF2
        text = ""
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    except ImportError:
        return "Error: PyPDF2 not installed. Run: pip install PyPDF2"
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"


def extract_text_from_docx(file_path):
    """Extract text from DOCX file"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return '\n'.join(text)
    except ImportError:
        return "Error: python-docx not installed. Run: pip install python-docx"
    except Exception as e:
        return f"Error extracting DOCX: {str(e)}"


def extract_text_from_doc(file_path):
    """Extract text from DOC file (older format)"""
    try:
        import textract
        text = textract.process(file_path).decode('utf-8')
        return text
    except ImportError:
        return "Error: textract not installed. Run: pip install textract"
    except Exception as e:
        # Fallback: try to read as plain text
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except:
            return f"Error extracting DOC: {str(e)}"


def extract_text_from_pptx(file_path):
    """Extract text from PPTX file"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"


def extract_text_from_ppt(file_path):
    """Extract text from PPT file (older format)"""
    try:
        import textract
        text = textract.process(file_path).decode('utf-8')
        return text
    except ImportError:
        return "Error: textract not installed. Run: pip install textract"
    except Exception as e:
        return f"Error extracting PPT: {str(e)}"


def extract_text_from_file(file_path):
    """
    Main function to extract text from any supported file format.
    Auto-detects format based on file extension.
    """
    if not os.path.exists(file_path):
        return "Error: File not found"
    
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    extractors = {
        '.txt': extract_text_from_txt,
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.doc': extract_text_from_doc,
        '.pptx': extract_text_from_pptx,
        '.ppt': extract_text_from_ppt,
    }
    
    extractor = extractors.get(ext)
    if not extractor:
        return f"Error: Unsupported file format '{ext}'. Supported: {', '.join(extractors.keys())}"
    
    return extractor(file_path)


def parse_question_from_text(text):
    """
    Parse question details from extracted text.
    Expected format:
    Title: <question title>
    Difficulty: <Easy/Medium/Hard>
    Description: <full description>
    
    Or just plain text which will be used as description.
    """
    lines = text.strip().split('\n')
    
    title = None
    difficulty = 'Medium'
    description = text
    
    # Try to parse structured format
    title_match = re.search(r'(?:Title|Problem):\s*(.+)', text, re.IGNORECASE)
    if title_match:
        title = title_match.group(1).strip()
    
    diff_match = re.search(r'Difficulty:\s*(Easy|Medium|Hard)', text, re.IGNORECASE)
    if diff_match:
        difficulty = diff_match.group(1).capitalize()
    
    desc_match = re.search(r'(?:Description|Problem Statement):\s*(.+)', text, re.IGNORECASE | re.DOTALL)
    if desc_match:
        description = desc_match.group(1).strip()
    
    # If no title found, use first line or first 50 chars
    if not title:
        first_line = lines[0].strip() if lines else "Untitled Problem"
        title = first_line[:100] if len(first_line) <= 100 else first_line[:97] + "..."
    
    return {
        'title': title,
        'difficulty': difficulty,
        'description': description
    }


if __name__ == '__main__':
    # Test the extractor
    import sys
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        text = extract_text_from_file(file_path)
        print("Extracted Text:")
        print("=" * 50)
        print(text)
        print("=" * 50)
        
        question = parse_question_from_text(text)
        print("\nParsed Question:")
        print(f"Title: {question['title']}")
        print(f"Difficulty: {question['difficulty']}")
        print(f"Description: {question['description'][:200]}...")
    else:
        print("Usage: python file_extractor.py <file_path>")
